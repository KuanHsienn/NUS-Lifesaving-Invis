import gc
import glob
import os
import platform
import shutil
import subprocess
import time

import pandas as pd
from pptx import Presentation


def replace_placeholders(prs, name):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{name}}" in run.text:
                            run.text = run.text.replace("{{name}}", name)


def _find_libreoffice():
    for cmd in ["libreoffice", "soffice"]:
        if shutil.which(cmd):
            return cmd
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if platform.system() == "Darwin" and os.path.exists(mac_path):
        return mac_path
    return None


def _convert_via_libreoffice(lo_cmd, pptx_path, pdf_path):
    out_dir = os.path.dirname(os.path.abspath(pdf_path))
    subprocess.run(
        [lo_cmd, "--headless", "--convert-to", "pdf", "--outdir", out_dir,
         os.path.abspath(pptx_path)],
        capture_output=True, check=True,
    )
    # LibreOffice names the output after the input stem
    generated = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    if os.path.abspath(generated) != os.path.abspath(pdf_path):
        os.rename(generated, pdf_path)


def _iter_names(df):
    for index, row in df.iterrows():
        if index == 0:
            continue
        for name in str(row["Name"]).replace("\n", ",").split(","):
            name = name.strip()
            if name:
                yield name


def _find_powerpoint_mac():
    return os.path.exists("/Applications/Microsoft PowerPoint.app")


def _convert_via_powerpoint_mac(pptx_path, pdf_path):
    abs_pptx = os.path.abspath(pptx_path).replace('"', '\\"')
    abs_pdf = os.path.abspath(pdf_path).replace('"', '\\"')
    script = (
        f'tell application "Microsoft PowerPoint"\n'
        f'  set thePres to open POSIX file "{abs_pptx}"\n'
        f'  save thePres in POSIX file "{abs_pdf}" as save as PDF\n'
        f'  close thePres saving no\n'
        f'end tell'
    )
    subprocess.run(["osascript", "-e", script], check=True, capture_output=True)


def _generate_macos(df, template_path, cert_output_dir):
    if _find_powerpoint_mac():
        print("Using PowerPoint for Mac for PDF export...")
        converter = "powerpoint"
    else:
        lo_cmd = _find_libreoffice()
        if lo_cmd is None:
            print(
                "Neither PowerPoint for Mac nor LibreOffice found.\n"
                "Install Microsoft PowerPoint or LibreOffice to generate PDFs on macOS."
            )
            return
        print(f"Using LibreOffice for PDF export...")
        converter = lo_cmd

    for single_name in _iter_names(df):
        safe_name = single_name.replace(" ", "_").upper()
        pptx_path = os.path.join(cert_output_dir, f"{safe_name}_CERT.pptx")
        pdf_path  = os.path.join(cert_output_dir, f"{safe_name}_CERT.pdf")

        if os.path.exists(pdf_path):
            if os.path.exists(pptx_path):
                os.remove(pptx_path)
            continue

        prs = Presentation(template_path)
        replace_placeholders(prs, single_name)
        prs.save(pptx_path)

        try:
            if converter == "powerpoint":
                _convert_via_powerpoint_mac(pptx_path, pdf_path)
            else:
                _convert_via_libreoffice(converter, pptx_path, pdf_path)
            print(f"Generated PDF for: {single_name}")
        except subprocess.CalledProcessError as e:
            print(f"  Warning: PDF conversion failed for {single_name} — {e.stderr}")
        finally:
            if os.path.exists(pptx_path):
                os.remove(pptx_path)

    print("Certificate generation complete.")


def generate_certificates(script_dir, output_dir):
    cert_output_dir = os.path.join(output_dir, "Certificates")
    os.makedirs(cert_output_dir, exist_ok=True)

    template_path = os.path.join(script_dir, "certificate_template.pptx")
    if not os.path.exists(template_path):
        print(f"Template not found at {template_path}. Skipping Certs.")
        return

    matches = [f for f in glob.glob(os.path.join(script_dir, "*.xlsx"))
               if "volunteers" in os.path.basename(f).lower()]
    if not matches:
        print("No 'volunteers.xlsx' found for certificate generation.")
        return

    df = pd.read_excel(matches[0], header=None, names=["Name"])
    df["Name"] = df["Name"].str.strip()
    df = df.dropna(subset=["Name"]).drop_duplicates(subset="Name")

    system = platform.system()
    if system == "Windows":
        _generate_windows(df, template_path, cert_output_dir)
    elif system == "Darwin":
        _generate_macos(df, template_path, cert_output_dir)
    else:
        lo_cmd = _find_libreoffice()
        if lo_cmd is None:
            print("LibreOffice not found. Install LibreOffice to generate PDFs on this platform.")
            return
        _generate_libreoffice(df, template_path, cert_output_dir, lo_cmd)


def _generate_windows(df, template_path, cert_output_dir):
    import win32com.client

    print("Opening PowerPoint for PDF export...")
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    try:
        for single_name in _iter_names(df):
            safe_name = single_name.replace(" ", "_").upper()
            pptx_path = os.path.join(cert_output_dir, f"{safe_name}_CERT.pptx")
            pdf_path  = os.path.join(cert_output_dir, f"{safe_name}_CERT.pdf")

            if os.path.exists(pdf_path):
                if os.path.exists(pptx_path):
                    os.remove(pptx_path)
                continue

            prs = Presentation(template_path)
            replace_placeholders(prs, single_name)
            prs.save(pptx_path)

            try:
                pres = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
                pres.SaveAs(os.path.abspath(pdf_path), 32)
                pres.Close()
                del pres
                gc.collect()
                time.sleep(0.5)
                print(f"Generated PDF for: {single_name}")
            finally:
                if os.path.exists(pptx_path):
                    os.remove(pptx_path)
    finally:
        powerpoint.Quit()
        print("Certificate generation complete.")


def _generate_libreoffice(df, template_path, cert_output_dir, lo_cmd):
    print(f"Using LibreOffice for PDF export...")
    for single_name in _iter_names(df):
        safe_name = single_name.replace(" ", "_").upper()
        pptx_path = os.path.join(cert_output_dir, f"{safe_name}_CERT.pptx")
        pdf_path  = os.path.join(cert_output_dir, f"{safe_name}_CERT.pdf")

        if os.path.exists(pdf_path):
            if os.path.exists(pptx_path):
                os.remove(pptx_path)
            continue

        prs = Presentation(template_path)
        replace_placeholders(prs, single_name)
        prs.save(pptx_path)

        try:
            _convert_via_libreoffice(lo_cmd, pptx_path, pdf_path)
            print(f"Generated PDF for: {single_name}")
        except subprocess.CalledProcessError as e:
            print(f"  Warning: PDF conversion failed for {single_name} — {e.stderr}")
        finally:
            if os.path.exists(pptx_path):
                os.remove(pptx_path)

    print("Certificate generation complete.")
